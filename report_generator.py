import os
from pptx import Presentation
from pptx.util import Inches

IMAGES_DIR = os.path.join(os.path.dirname(__file__), 'images')

def collect_images():
    if not os.path.exists(IMAGES_DIR):
        return []
    imgs = []
    for f in sorted(os.listdir(IMAGES_DIR)):
        if f.lower().endswith(('.png', '.jpg', '.jpeg')):
            imgs.append(os.path.join(IMAGES_DIR, f))
    return imgs


def add_title_slide(prs, title, subtitle=''):
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = title
    if subtitle:
        slide.placeholders[1].text = subtitle


def add_image_slide(prs, image_path, title=None):
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    left = top = Inches(0.5)
    pic = slide.shapes.add_picture(image_path, left, top, width=Inches(9))
    if title:
        # add a small title textbox
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(9), Inches(0.5))
        tx.text = title


def main(output='report.pptx'):
    imgs = collect_images()
    if not imgs:
        print('No images found in', IMAGES_DIR)
        return

    prs = Presentation()
    add_title_slide(prs, 'Automated Report', 'Generated from images/ outputs')

    for img in imgs:
        name = os.path.basename(img)
        add_image_slide(prs, img, title=name)

    out_path = os.path.join(os.path.dirname(__file__), output)
    prs.save(out_path)
    print('Saved PPTX report to', out_path)


if __name__ == '__main__':
    main()
